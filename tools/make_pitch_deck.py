from pptx import Presentation
from pptx.util import Pt

slides = [
    {
        "title": "Problem",
        "bullets": [
            "Who: Field technicians and ops teams at mid-size manufacturing and logistics companies.",
            "Pain: They lose hours daily to manual troubleshooting, delayed diagnostics, and unclear root causes.",
            "Impact: Downtime and slow fixes cost an estimated 5–10% of operating time and revenue.",
            "Example: A single machine outage can cascade across production lines, delaying shipments and increasing costs.",
        ],
    },
    {
        "title": "Proposed solution & key features",
        "bullets": [
            "One-line: Lightweight automated diagnostics and incident-resolution assistant that reduces mean-time-to-repair.",
            "Real-time diagnostics and anomaly detection.",
            "Step-by-step guided repair actions for technicians.",
            "Incident logging, root-cause hints, and suggested fixes.",
            "Integration with existing monitoring and ticketing systems.",
        ],
    },
    {
        "title": "Tools / Tech stack",
        "bullets": [
            "Frontend: React or lightweight web UI.",
            "Backend: FastAPI (Python) or Node/Express.",
            "ML/Analytics: PyTorch / scikit-learn; optional LLMs (OpenAI/local).",
            "Infra: Docker + CI/CD; AWS/GCP/Azure or on-prem.",
            "Storage & Telemetry: PostgreSQL, Redis, MQTT/ROS2.",
        ],
    },
    {
        "title": "ICP (Target customer)",
        "bullets": [
            "Industry: Manufacturing, warehousing, logistics, industrial automation.",
            "Buyer personas: Operations manager, maintenance lead, reliability engineer.",
            "Early adopters: Companies with 50–500 employees and distributed shop floors.",
            "Why they'll pay: Faster mean-time-to-repair yields measurable ROI through increased uptime and lower labor cost.",
        ],
    },
]


def make_pptx(path="phase1_pitch_deck.pptx"):
    prs = Presentation()
    for s in slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = s["title"]
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        if s["bullets"]:
            body.text = s["bullets"][0]
            for b in s["bullets"][1:]:
                p = body.add_paragraph()
                p.text = b
                p.level = 0
    prs.save(path)


if __name__ == "__main__":
    import os

    out = os.path.join(os.getcwd(), "phase1_pitch_deck.pptx")
    make_pptx(out)
    print(f"Saved: {out}")
